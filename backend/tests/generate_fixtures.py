import os
import wave
import json
import struct
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent / "fixtures"
FIXTURES_DIR.mkdir(parents=True, exist_ok=True)

def generate_fixtures():
    print(f"Generating test fixtures in {FIXTURES_DIR}...")

    # 1. TXT fixture
    txt_path = FIXTURES_DIR / "sample.txt"
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write("FILE CONV Test Document\n\nThis is a sample plain text document for conversion testing.")

    # 2. MD fixture
    md_path = FIXTURES_DIR / "sample.md"
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("# FILE CONV Sample Markdown\n\n- Feature 1\n- Feature 2\n\n**Bold Text** and *Italic Text*")

    # 3. CSV fixture
    csv_path = FIXTURES_DIR / "sample.csv"
    with open(csv_path, "w", encoding="utf-8") as f:
        f.write("name,role,department\nAlice,Developer,Engineering\nBob,Designer,UX\nCharlie,Manager,Product\n")

    # 4. JSON fixture
    json_path = FIXTURES_DIR / "sample.json"
    data = [
        {"id": 1, "product": "Widget A", "price": 9.99},
        {"id": 2, "product": "Widget B", "price": 19.99}
    ]
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)

    # 5. HTML fixture
    html_path = FIXTURES_DIR / "sample.html"
    with open(html_path, "w", encoding="utf-8") as f:
        f.write("<!DOCTYPE html><html><head><title>Test HTML</title></head><body><h1>FILE CONV Web Page</h1><p>Sample HTML paragraph for batch conversion test.</p><ul><li>Item A</li><li>Item B</li></ul></body></html>")

    # 6. XLSX fixture (via openpyxl)
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["Name", "Score", "City"])
        ws.append(["Diana", 95, "New York"])
        ws.append(["Edward", 88, "London"])
        wb.save(FIXTURES_DIR / "sample.xlsx")
    except Exception as e:
        print(f"Failed creating XLSX fixture: {e}")

    # 7. PNG fixture (via Pillow)
    try:
        from PIL import Image, ImageDraw
        img = Image.new("RGBA", (200, 200), color=(73, 109, 137, 255))
        d = ImageDraw.Draw(img)
        d.text((10, 80), "FILE CONV", fill=(255, 255, 0, 255))
        img.save(FIXTURES_DIR / "sample.png", format="PNG")
        img.convert("RGB").save(FIXTURES_DIR / "sample.jpg", format="JPEG")
    except Exception as e:
        print(f"Failed creating PNG/JPG fixture: {e}")

    # 8. DOCX fixture (via python-docx)
    try:
        import docx
        doc = docx.Document()
        doc.add_heading("FILE CONV Test DOCX", 0)
        doc.add_paragraph("This is a sample word document paragraph.")
        doc.save(FIXTURES_DIR / "sample.docx")

        doc2 = docx.Document()
        doc2.add_heading("FILE CONV Second DOCX", 0)
        doc2.add_paragraph("This is the second document for merge testing.")
        doc2.save(FIXTURES_DIR / "sample2.docx")
    except Exception as e:
        print(f"Failed creating DOCX fixture: {e}")

    # 9. PDF fixture (via PyMuPDF fitz)
    try:
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "FILE CONV Test PDF Page 1", fontsize=18)
        doc.save(FIXTURES_DIR / "sample.pdf")
        doc.close()

        doc2 = fitz.open()
        page2 = doc2.new_page()
        page2.insert_text((50, 50), "FILE CONV Test PDF Document 2 Page 1", fontsize=18)
        doc2.save(FIXTURES_DIR / "sample2.pdf")
        doc2.close()
    except Exception as e:
        print(f"Failed creating PDF fixture: {e}")

    # 10. PPTX fixture (via python-pptx)
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "FILE CONV PPTX Test"
        slide.placeholders[1].text = "Sample presentation deck."
        prs.save(FIXTURES_DIR / "sample.pptx")

        prs2 = Presentation()
        slide2 = prs2.slides.add_slide(prs2.slide_layouts[1])
        slide2.shapes.title.text = "Second PPTX Deck"
        slide2.placeholders[1].text = "Slide content for deck 2."
        prs2.save(FIXTURES_DIR / "sample2.pptx")
    except Exception as e:
        print(f"Failed creating PPTX fixture: {e}")

    # 11. WAV fixture
    try:
        wav_path = FIXTURES_DIR / "sample.wav"
        with wave.open(str(wav_path), "wb") as wav_file:
            wav_file.setnchannels(1)
            wav_file.setsampwidth(2)
            wav_file.setframerate(44100)
            import math
            frames = []
            for i in range(44100):
                val = int(32767.0 * 0.3 * math.sin(2.0 * math.pi * 440.0 * i / 44100.0))
                frames.append(struct.pack("<h", val))
            wav_file.writeframes(b"".join(frames))
    except Exception as e:
        print(f"Failed creating WAV fixture: {e}")

    # 12. MP3 & MP4 fixture
    try:
        import subprocess
        ffmpeg_bin = os.path.expanduser(r"~\AppData\Local\Microsoft\WinGet\Packages\Gyan.FFmpeg.Essentials_Microsoft.Winget.Source_8wekyb3d8bbwe\ffmpeg-8.1.1-essentials_build\bin\ffmpeg.exe")
        if not os.path.exists(ffmpeg_bin):
            ffmpeg_bin = "ffmpeg"
        
        subprocess.run([ffmpeg_bin, "-y", "-i", str(FIXTURES_DIR / "sample.wav"), str(FIXTURES_DIR / "sample.mp3")], capture_output=True)
        subprocess.run([ffmpeg_bin, "-y", "-f", "lavfi", "-i", "color=c=blue:s=320x240:d=1", "-f", "lavfi", "-i", "sine=f=440:d=1", str(FIXTURES_DIR / "sample.mp4")], capture_output=True)
    except Exception as e:
        print(f"Failed creating MP3/MP4 fixtures: {e}")

    print("Fixtures generation complete!")

if __name__ == "__main__":
    generate_fixtures()
