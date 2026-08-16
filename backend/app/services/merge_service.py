import os
import logging
from typing import List, Tuple, Optional
from app.core.config import settings

logger = logging.getLogger("merge_service")

class MergeService:
    @staticmethod
    def merge_pdf_files(file_paths: List[str], output_path: str) -> Tuple[bool, Optional[str]]:
        """Merge multiple PDF files into one output PDF using PyMuPDF."""
        try:
            import fitz
            merged_doc = fitz.open()

            for path in file_paths:
                if not os.path.exists(path):
                    return False, f"File not found: {os.path.basename(path)}"
                doc = fitz.open(path)
                merged_doc.insert_pdf(doc)
                doc.close()

            merged_doc.save(output_path)
            merged_doc.close()
            return True, None
        except Exception as e:
            return False, f"PDF Merge error: {str(e)}"

    @staticmethod
    def merge_pptx_files(file_paths: List[str], output_path: str) -> Tuple[bool, Optional[str]]:
        """Merge multiple PPTX files into one presentation deck using python-pptx."""
        try:
            from pptx import Presentation
            
            if not file_paths:
                return False, "No input PPTX files provided."

            base_prs = Presentation(file_paths[0])

            for path in file_paths[1:]:
                if not os.path.exists(path):
                    return False, f"File not found: {os.path.basename(path)}"
                prs = Presentation(path)
                for slide in prs.slides:
                    # Determine slide layout matching slide master
                    slide_layout = base_prs.slide_layouts[6]  # Blank layout
                    new_slide = base_prs.slides.add_slide(slide_layout)
                    
                    # Copy shapes & text from source slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            txBox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                            tf = txBox.text_frame
                            tf.text = shape.text

            base_prs.save(output_path)
            return True, None
        except Exception as e:
            return False, f"PPTX Merge error: {str(e)}"

    @staticmethod
    def merge_docx_files(file_paths: List[str], output_path: str) -> Tuple[bool, Optional[str]]:
        """Merge multiple Word DOCX files into one document using python-docx."""
        try:
            import docx

            if not file_paths:
                return False, "No input DOCX files provided."

            base_doc = docx.Document(file_paths[0])

            for path in file_paths[1:]:
                if not os.path.exists(path):
                    return False, f"File not found: {os.path.basename(path)}"
                doc = docx.Document(path)
                
                # Add page break before appending next document
                base_doc.add_page_break()

                # Append element bodies
                for element in doc.element.body:
                    base_doc.element.body.append(element)

            base_doc.save(output_path)
            return True, None
        except Exception as e:
            return False, f"DOCX Merge error: {str(e)}"
