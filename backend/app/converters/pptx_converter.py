import os
import asyncio
import logging
from typing import Dict, Any, Callable, Optional, Tuple
from app.converters.base import BaseConverter
from app.core.config import settings

logger = logging.getLogger("pptx_converter")

class PPTXConverter(BaseConverter):
    async def convert(
        self,
        input_path: str,
        output_path: str,
        source_ext: str,
        target_ext: str,
        options: Optional[Dict[str, Any]] = None,
        progress_callback: Optional[Callable[[int, str], None]] = None
    ) -> Tuple[bool, Optional[str]]:
        options = options or {}
        target_ext = target_ext.lower().lstrip(".")

        try:
            if progress_callback:
                progress_callback(10, f"Loading PowerPoint presentation ({source_ext.upper()})...")

            # PPT / PPTX -> Markdown (.md)
            if target_ext == "md":
                return self._convert_to_markdown(input_path, output_path, progress_callback)

            # PPT / PPTX -> PDF
            elif target_ext == "pdf":
                return await self._convert_to_pdf(input_path, output_path, progress_callback)

            # PPT / PPTX -> Image (PNG/JPG)
            elif target_ext in ["png", "jpg", "jpeg"]:
                # Convert PPTX to PDF first, then render PDF Page 1 to PNG/JPG via PyMuPDF
                temp_pdf = input_path + ".temp.pdf"
                success, err = await self._convert_to_pdf(input_path, temp_pdf, progress_callback)
                if not success or not os.path.exists(temp_pdf):
                    return False, f"Failed converting presentation to PDF for image rendering: {err}"

                try:
                    import fitz
                    doc = fitz.open(temp_pdf)
                    if len(doc) > 0:
                        page = doc[0]
                        pix = page.get_pixmap(dpi=150)
                        pix.save(output_path)
                    doc.close()
                    os.remove(temp_pdf)
                    if progress_callback:
                        progress_callback(100, "Presentation slide rendered as image.")
                    return True, None
                except Exception as e:
                    if os.path.exists(temp_pdf):
                        os.remove(temp_pdf)
                    return False, f"Slide image rendering error: {str(e)}"

            else:
                return False, f"Unsupported PPTX target format: {target_ext}"

        except Exception as e:
            return False, f"PowerPoint conversion error: {str(e)}"

    def _convert_to_markdown(self, input_path: str, output_path: str, progress_callback: Optional[Callable] = None) -> Tuple[bool, Optional[str]]:
        try:
            from pptx import Presentation
        except ImportError:
            return False, "python-pptx library is not installed."

        try:
            prs = Presentation(input_path)
            md_lines = []
            total_slides = len(prs.slides)

            md_lines.append("# Presentation Overview\n")

            for idx, slide in enumerate(prs.slides):
                md_lines.append(f"## Slide {idx + 1}\n")
                
                # Extract text shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder and shape.placeholder_format.type == 1:
                            md_lines.append(f"### {text}\n")
                        else:
                            for line in text.splitlines():
                                if line.strip():
                                    md_lines.append(f"- {line.strip()}")
                            md_lines.append("")

                # Extract speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        md_lines.append(f"\n> **Speaker Notes:** {notes}\n")

                md_lines.append("\n---\n")

                if progress_callback and total_slides > 0:
                    pct = int(10 + (idx + 1) / total_slides * 80)
                    progress_callback(pct, f"Extracting slide {idx + 1}/{total_slides}")

            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n".join(md_lines))

            if progress_callback:
                progress_callback(100, "PPTX converted to Markdown.")
            return True, None

        except Exception as e:
            return False, f"Failed parsing PPTX to Markdown: {str(e)}"

    async def _convert_to_pdf(self, input_path: str, output_path: str, progress_callback: Optional[Callable] = None) -> Tuple[bool, Optional[str]]:
        abs_input = os.path.abspath(input_path)
        abs_output = os.path.abspath(output_path)

        # 1. Try Windows Native PowerPoint COM automation if on Windows
        if os.name == "nt":
            try:
                import win32com.client
                if progress_callback:
                    progress_callback(30, "Exporting presentation via PowerPoint COM engine...")
                
                # Run COM in thread pool to prevent blocking loop
                def export_com():
                    import pythoncom
                    pythoncom.CoInitialize()
                    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                    # PpSaveAsFileType.ppSaveAsPDF = 32
                    deck = ppt_app.Presentations.Open(abs_input, WithWindow=False)
                    deck.SaveAs(abs_output, 32)
                    deck.Close()
                    ppt_app.Quit()

                await asyncio.to_thread(export_com)

                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    if progress_callback:
                        progress_callback(100, "PowerPoint PDF export complete.")
                    return True, None
            except Exception as e:
                logger.warning(f"PowerPoint COM export failed, falling back to PDF text layout builder: {e}")

        # 2. Fallback: Build structured PDF deck using python-pptx + PyMuPDF
        try:
            from pptx import Presentation
            import fitz

            if progress_callback:
                progress_callback(40, "Building PDF layout from presentation slides...")

            prs = Presentation(input_path)
            pdf_doc = fitz.open()

            for idx, slide in enumerate(prs.slides):
                # Create 16:9 widescreen landscape page (960x540 points)
                page = pdf_doc.new_page(width=960, height=540)
                
                # Add Header / Slide Title background
                page.draw_rect(fitz.Rect(0, 0, 960, 60), color=(0.15, 0.25, 0.45), fill=(0.15, 0.25, 0.45))
                page.insert_text((30, 40), f"Slide {idx + 1}", fontsize=20, color=(1, 1, 1))

                y_pos = 100
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        lines = text.splitlines()
                        for line in lines:
                            if y_pos > 500:
                                break
                            page.insert_text((40, y_pos), line[:110], fontsize=12, color=(0.1, 0.1, 0.1))
                            y_pos += 24

            pdf_doc.save(output_path)
            pdf_doc.close()

            if progress_callback:
                progress_callback(100, "PPTX to PDF conversion complete.")
            return True, None

        except Exception as e:
            return False, f"PPTX to PDF export failed: {str(e)}"
